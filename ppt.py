import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import openai

# Set your OpenAI API key
openai.api_type = "azure"
openai.api_version = "2023-07-01-preview"
openai.api_base = "https://hsbcoai.openai.azure.com/"
openai.api_key = "324de9c5fb0444278d98f92dbee4ee2a"

def generate_content(prompt):
    # Use OpenAI to generate content based on the prompt
    response = openai.Completion.create(
        engine="gpt-35-turbo",
        prompt=prompt,
        max_tokens=140,  # Adjust as needed
        temperature=0.2,  # Adjust for creativity
    )
    return response.choices[0].text

def create_presentation(topic):
    prs = Presentation('template.pptx')
    
    slide_titles = ["Introduction", "Main Content", "Conclusion"]
    slide_layouts = prs.slide_layouts

    for title in slide_titles:
        slide_layout = slide_layouts[1]  # Use Title Slide layout (index 1) for other slides
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Change the font type of titles to something else and make them bold
        title_text_frame = title_placeholder.text_frame
        title_text_frame.paragraphs[0].font.name = 'Palace Script MT'  # Change to your desired font
        title_text_frame.paragraphs[0].font.bold = True

        # Generate content for the content placeholder using OpenAI
        content_prompt = f"Write 4 key points about the '{topic}' for '{title}' slide. The points should only be related to the  '{topic}' and no irrelevant information should be added:"
        content = generate_content(content_prompt)

        # Add bullet points with a maximum of 4 points per slide
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        paragraphs = content.split('\n')
        max_points = 4
        points_added = 0

        for paragraph in paragraphs:
            if paragraph.strip():
                p = text_frame.add_paragraph()
                p.text = paragraph
                p.space_after = Pt(12)
                p.space_before = Pt(0)
                p.level = 0

                # Set the font size to 18 Pt and change to Times New Roman
                for run in p.runs:
                    font = run.font
                    font.size = Pt(10)
                    font.name = 'Times New Roman'

                points_added += 1

                if points_added >= max_points:
                    break

    # Add the HSBC logo to all slides
    logo_path = 'hsbc.png'
    left = Inches(8.2)  # Adjust the left position as needed
    top = Inches(0.5)  # Adjust the top position as needed
    for slide in prs.slides:
        pic = slide.shapes.add_picture(logo_path, left, top, height=Inches(1))  # Adjust the height as needed

    # Replace the topic on the first slide
    if prs.slides:
        first_slide = prs.slides[0]
        title_placeholder = first_slide.shapes.title
        title_placeholder.text = topic

    return prs

def main():
    st.title("PowerPoint Presentation Generator")

    # Get user input for the presentation topic
    topic = st.text_input("Enter the topic for the presentation:")

    if st.button("Generate Presentation"):
        if topic:
            st.text("Generating the presentation. Please wait...")

            # Create the presentation
            presentation = create_presentation(topic)

            # Save the presentation as a new file
            presentation_file_name = f"{topic}_presentation.pptx"
            presentation.save(presentation_file_name)

            st.success(f"Presentation generated and saved as {presentation_file_name}")

if __name__ == "__main__":
    main()

