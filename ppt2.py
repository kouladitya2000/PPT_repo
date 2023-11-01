import streamlit as st
from pptx import Presentation
from pptx.util import Pt
import openai

# Set your OpenAI API key
openai.api_type = "azure"
openai.api_version = "2023-07-01-preview"
openai.api_base = "https://hsbcoai.openai.azure.com/"
openai.api_key = "324de9c5fb0444278d98f92dbee4ee2a"

def generate_content(prompt):
    # Use OpenAI to generate content based on the prompt
    response = openai.Completion.create(
        engine="gpt-35-turbo",
        prompt=prompt,
        max_tokens=140,  # Adjust as needed
        temperature=0.2,  # Adjust for creativity
    )
    return response.choices[0].text

def create_presentation(topic):
    prs = Presentation()
    slide_titles = ["Introduction", "Main Content", "Conclusion"]
    slide_layouts = prs.slide_layouts

    for title in slide_titles:
        slide_layout = slide_layouts[1]  # Use Title Slide layout (index 1) for all slides
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Generate content for the content placeholder using OpenAI
        content_prompt = f"Write 4 key points about the '{topic}' for '{title}' slide. The points should only be related to the  '{topic}' and no irrelevant information should be added:"
        content = generate_content(content_prompt)

        # Add bullet points with a maximum of 4 points per slide
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        paragraphs = content.split('\n')
        max_points = 4
        points_added = 0

        for paragraph in paragraphs:
            if paragraph.strip():
                p = text_frame.add_paragraph()
                p.text = paragraph
                p.space_after = Pt(12)
                p.space_before = Pt(0)
                p.level = 0

                # Set the font size to 18 Pt
                for run in p.runs:
                    font = run.font
                    font.size = Pt(10)

                points_added += 1

                if points_added >= max_points:
                    break

    return prs

def main():
    st.title("PowerPoint Presentation Generator")

    # Get user input for the presentation topic
    topic = st.text_input("Enter the topic for the presentation:")

    if st.button("Generate Presentation"):
        if topic:
            st.text("Generating the presentation. Please wait...")

            # Create the presentation
            presentation = create_presentation(topic)

            # Save the presentation as a file
            presentation_file_name = f"{topic}_presentation.pptx"
            presentation.save(presentation_file_name)

            st.success(f"Presentation generated and saved as {presentation_file_name}")

if __name__ == "__main__":
    main()
